import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create Presentation
prs = Presentation()
# Set Widescreen (16:9)
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Styling Constants
PRIMARY_DARK = RGBColor(26, 54, 93)   # #1a365d (Deep Corporate Blue)
ACCENT_BLUE = RGBColor(43, 108, 176)   # #2b6cb0 (Accent Blue)
TEXT_DARK = RGBColor(45, 55, 72)     # #2d3748 (Charcoal Text)
TEXT_MUTED = RGBColor(113, 128, 150)  # #718096 (Grey Text)
WHITE = RGBColor(255, 255, 255)
BG_LIGHT = RGBColor(247, 250, 252)    # #f7fafc (Light Grey)

# Helper function to add a standard slide header
def add_slide_header(slide, title_text):
    # Title Text Box
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Georgia'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_DARK
    
    # Horizontal separator line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.15), Inches(11.83), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.color.rgb = ACCENT_BLUE

# Helper function to set formatting for paragraphs
def format_bullet(p, text, is_bold_prefix=True, font_size=14):
    p.font.name = 'Calibri'
    p.font.size = Pt(font_size)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(8)
    
    if is_bold_prefix and ":" in text:
        parts = text.split(":", 1)
        p.text = ""
        run1 = p.add_run()
        run1.text = parts[0] + ":"
        run1.font.bold = True
        run1.font.color.rgb = PRIMARY_DARK
        
        run2 = p.add_run()
        run2.text = parts[1]
        run2.font.bold = False
    else:
        p.text = text

# Helper function to insert an image with border/caption
def insert_chart(slide, image_path, left, top, width, height, caption=""):
    if os.path.exists(image_path):
        # Add a subtle border shape behind the chart
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left - Inches(0.05), top - Inches(0.05), width + Inches(0.1), height + Inches(0.1))
        border.fill.solid()
        border.fill.fore_color.rgb = WHITE
        border.line.color.rgb = RGBColor(226, 232, 240) # Light grey border
        
        # Add the image
        slide.shapes.add_picture(image_path, left, top, width, height)
        
        if caption:
            caption_box = slide.shapes.add_textbox(left, top + height + Inches(0.05), width, Inches(0.3))
            ctf = caption_box.text_frame
            ctf.word_wrap = True
            cp = ctf.paragraphs[0]
            cp.text = caption
            cp.alignment = PP_ALIGN.CENTER
            cp.font.name = 'Calibri'
            cp.font.size = Pt(10)
            cp.font.italic = True
            cp.font.color.rgb = TEXT_MUTED
    else:
        # Placeholder shape
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(241, 245, 249)
        placeholder.line.color.rgb = ACCENT_BLUE
        tf = placeholder.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[Chart Not Found: {os.path.basename(image_path)}]"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Calibri'
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED

# ==========================================
# SLIDE 1: Title Slide (Bespoke design)
# ==========================================
slide_layout = prs.slide_layouts[6] # Blank
slide = prs.slides.add_slide(slide_layout)

# Add a dark blue decorative block on the left
left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.0), Inches(7.5))
left_bar.fill.solid()
left_bar.fill.fore_color.rgb = PRIMARY_DARK
left_bar.line.fill.background()

# Title text frame in the blue block
title_box_left = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(3.0), Inches(3.5))
ltf = title_box_left.text_frame
ltf.word_wrap = True
lp = ltf.paragraphs[0]
lp.text = "MUTUAL\nFUND\nANALYTICS"
lp.font.name = 'Georgia'
lp.font.size = Pt(36)
lp.font.bold = True
lp.font.color.rgb = WHITE
lp.space_after = Pt(20)

lp2 = ltf.add_paragraph()
lp2.text = "Capstone Project\nPresentation"
lp2.font.name = 'Calibri'
lp2.font.size = Pt(18)
lp2.font.color.rgb = ACCENT_BLUE

# Right side content block
right_box = slide.shapes.add_textbox(Inches(4.8), Inches(2.0), Inches(7.5), Inches(4.0))
rtf = right_box.text_frame
rtf.word_wrap = True

rp1 = rtf.paragraphs[0]
rp1.text = "End-to-End Quantitative Pipeline & Business Intelligence Dashboard"
rp1.font.name = 'Georgia'
rp1.font.size = Pt(28)
rp1.font.bold = True
rp1.font.color.rgb = PRIMARY_DARK
rp1.space_after = Pt(24)

rp2 = rtf.add_paragraph()
rp2.text = "A production-grade data engineering and quantitative analytics system designed for the Indian Mutual Fund industry."
rp2.font.name = 'Calibri'
rp2.font.size = Pt(16)
rp2.font.color.rgb = TEXT_DARK
rp2.space_after = Pt(36)

rp3 = rtf.add_paragraph()
rp3.text = "Prepared for: BlueStock Intern Project Evaluation\nDeveloped by: Capstone Technical Team\nDate: June 2026"
rp3.font.name = 'Calibri'
rp3.font.size = Pt(13)
rp3.font.color.rgb = TEXT_MUTED

# Decorative divider line on the right side
div_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), Inches(1.8), Inches(7.5), Inches(0.04))
div_line.fill.solid()
div_line.fill.fore_color.rgb = ACCENT_BLUE
div_line.line.color.rgb = ACCENT_BLUE


# ==========================================
# SLIDE 2: Executive Summary
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_slide_header(slide, "Executive Summary")

# Left Column (Core Accomplishments)
left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(7.25), Inches(5.2))
ltf = left_box.text_frame
ltf.word_wrap = True

bullets = [
    "Automated ETL Pipeline: Standardized and cleaned 10 raw CSV datasets representing over 90,000 transaction, master, and index entries, resolving data anomalies programmatically.",
    "Relational SQLite Database: Populated a Star Schema SQLite database (bluestock_mf.db) optimizing performance and storage, enforcing structural integrity and primary/foreign key relationships.",
    "Quantitative Financial Engine: Computed returns, CAGR, Sharpe and Sortino ratios, CAPM Alpha/Beta, and Maximum Drawdowns across 40 schemes over a historical 4.5-year horizon.",
    "Advanced Strategy Modeling: Developed Value at Risk (VaR), Conditional VaR (CVaR), Herfindahl-Hirschman Index (HHI) for sector concentrations, and customer cohort retention tracking.",
    "Interactive Analytics Dashboard: Created a responsive 6-page dashboard featuring drill-through scorecards, state-wise flows, demographic insights, and a personalized recommendation engine."
]

for i, bullet in enumerate(bullets):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    format_bullet(p, bullet)

# Right Column (Key Database Highlights Callout)
callout_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.5), Inches(4.0), Inches(5.0))
callout_bg.fill.solid()
callout_bg.fill.fore_color.rgb = PRIMARY_DARK
callout_bg.line.fill.background()

callout_box = slide.shapes.add_textbox(Inches(8.7), Inches(1.7), Inches(3.6), Inches(4.6))
ctf = callout_box.text_frame
ctf.word_wrap = True

cp = ctf.paragraphs[0]
cp.text = "DATABASE HIGHLIGHTS"
cp.font.name = 'Georgia'
cp.font.size = Pt(16)
cp.font.bold = True
cp.font.color.rgb = WHITE
cp.alignment = PP_ALIGN.CENTER
cp.space_after = Pt(20)

stats = [
    "• 40 Mutual Fund Schemes",
    "• 10 Interlinked Tables",
    "• 90,000+ Total Records",
    "• 32,778 Investor Transactions",
    "• 5,000 Unique Customers",
    "• 1,043,664 Cr Total AUM",
    "• 3,521.58M INR Inflow Value"
]
for stat in stats:
    p = ctf.add_paragraph()
    p.text = stat
    p.font.name = 'Calibri'
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.space_after = Pt(10)


# ==========================================
# SLIDE 3: System Architecture
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_slide_header(slide, "System Architecture & Decoupled Data Pipeline")

# Description
intro_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.83), Inches(0.8))
itf = intro_box.text_frame
itf.word_wrap = True
ip = itf.paragraphs[0]
ip.text = "The system is structured as a modern data engineering pipeline, dividing operations into clear layers:"
ip.font.name = 'Calibri'
ip.font.size = Pt(15)
ip.font.color.rgb = TEXT_DARK

# Draw Architecture Blocks
block_width = Inches(2.6)
block_height = Inches(4.0)
top_pos = Inches(2.2)

layers = [
    {
        "num": "1",
        "title": "Raw Ingestion",
        "bg": RGBColor(237, 242, 249),
        "border": ACCENT_BLUE,
        "items": [
            "• 10 Ingested CSV Files",
            "• AMFI REST API Integration",
            "• Historical Daily NAVs",
            "• Investor Transactions",
            "• Portfolio Sector Weights",
            "• Index closing values"
        ]
    },
    {
        "num": "2",
        "title": "ETL & Storage",
        "bg": RGBColor(235, 248, 255),
        "border": ACCENT_BLUE,
        "items": [
            "• Python Standardization",
            "• Business-day Reindexing",
            "• Missing NAV Forward-fill",
            "• SQLite Normalized DB",
            "• Foreign Key Constraints",
            "• SQL indexing & queries"
        ]
    },
    {
        "num": "3",
        "title": "Quantitative Engine",
        "bg": RGBColor(224, 243, 255),
        "border": ACCENT_BLUE,
        "items": [
            "• Returns, CAGR, & MDD",
            "• Sharpe & Sortino ratios",
            "• CAPM Alpha & Beta Model",
            "• 95% Daily VaR & CVaR",
            "• HHI Concentration",
            "• Cohort retention matrices"
        ]
    },
    {
        "num": "4",
        "title": "Presentation Layer",
        "bg": PRIMARY_DARK,
        "border": PRIMARY_DARK,
        "items": [
            "• Python Server Backend",
            "• REST API endpoints",
            "• Vanilla JS Frontend",
            "• Dynamic Chart.js visuals",
            "• Multi-page filtering",
            "• Recommendation UI"
        ],
        "dark_mode": True
    }
]

for idx, layer in enumerate(layers):
    left_pos = Inches(0.75 + idx * 3.08)
    
    # Base shape
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, block_width, block_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = layer["bg"]
    shape.line.color.rgb = layer["border"]
    shape.line.width = Pt(1.5)
    
    # Text frame
    tb = slide.shapes.add_textbox(left_pos + Inches(0.1), top_pos + Inches(0.1), block_width - Inches(0.2), block_height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    # Layer Title
    tp = tf.paragraphs[0]
    tp.text = f"Layer {layer['num']}: {layer['title']}"
    tp.font.name = 'Georgia'
    tp.font.size = Pt(15)
    tp.font.bold = True
    tp.font.color.rgb = WHITE if layer.get("dark_mode") else PRIMARY_DARK
    tp.space_after = Pt(15)
    tp.alignment = PP_ALIGN.CENTER
    
    # Items
    for item in layer["items"]:
        p = tf.add_paragraph()
        p.text = item
        p.font.name = 'Calibri'
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE if layer.get("dark_mode") else TEXT_DARK
        p.space_after = Pt(8)

    # Arrow connector (except for the last block)
    if idx < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left_pos + block_width + Inches(0.1), top_pos + Inches(1.8), Inches(0.28), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_BLUE
        arrow.line.fill.background()


# ==========================================
# SLIDE 4: Database Schema Design
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_slide_header(slide, "Relational Database Star Schema")

# Left Column - Schema Layout
left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(5.0))
ltf = left_box.text_frame
ltf.word_wrap = True

bullets = [
    "dim_fund (Primary Dimension): Central master table representing 40 mutual fund schemes, mapping AMFI code, categories, benchmark, and expense ratios.",
    "fact_nav (Historical Dimension): Stores historical daily NAVs and returns reindexed to business days over 4.5 years (~46,000 rows).",
    "fact_performance (Analytical Dimension): Houses calculated composite scores, returns, Sharpe, Sortino, Alpha, Beta, MDD, and AUM.",
    "fact_transactions (Behavioral Fact): Tracks 32,778 transaction logs containing demographic variables, transaction values, and payment modes.",
    "fact_portfolio_holdings (Holding Fact): Maps portfolio stock symbols, weights, and sectors to each AMFI scheme."
]

for i, bullet in enumerate(bullets):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    format_bullet(p, bullet)

# Right Column - Auxiliary tables
right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.0))
rtf = right_box.text_frame
rtf.word_wrap = True

bullets_right = [
    "fact_advanced_risk (Risk Fact): Captures calculated 95% Value at Risk (VaR) and 95% Conditional VaR (CVaR) for all 40 schemes.",
    "fact_sector_concentration (Concentration Fact): Tracks Herfindahl-Hirschman Index (HHI) and concentration levels based on portfolio weightings.",
    "fact_sip_inflows (Industry Fact): Tracks aggregate monthly mutual fund SIP inflows, folio registrations, and AUM growth in India.",
    "fact_category_inflows (Industry Fact): Details net capital inflows across mutual fund categories on a monthly basis.",
    "dim_aum_fund_house & dim_industry_folio: Store aggregate industry parameters and fund house AUM statements."
]

for i, bullet in enumerate(bullets_right):
    p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
    format_bullet(p, bullet)


# ==========================================
# SLIDE 5: Data Pipeline & ETL Process
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_slide_header(slide, "Data Ingestion Pipeline & ETL Process")

# Main Box
main_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
mtf = main_box.text_frame
mtf.word_wrap = True

bullets = [
    "Holiday Reindexing & Forward-Filling: Daily NAVs in raw data (1,150 dates, Jan 2022 to May 2026) were missing holidays. The pipeline reindexes each scheme to standard business days (freq='B') and forward-fills (ffill) missing values. This prevents computation gaps and keeps chronological data consistent.",
    "Transaction Type Mapping & Normalization: Cleaned transactional metadata in investor logs. Standardized input strings to strict categories: SIP, Lumpsum, or Redemption, and discarded negative transaction amounts.",
    "KYC Status Standardization: Cleaned KYC fields, eliminating missing or corrupt strings and mapping them strictly to Verified or Pending to maintain compliance accuracy.",
    "Numeric Constraint Auditing: Verified that AUM, daily NAV, and transaction amount fields are strictly positive. Enforced expense ratios to standard ranges [0.1%, 2.5%] and converted drawdown percentages to negative values.",
    "Automated Quality Assurance: Created verify_pipeline.py which performs automated assertions to check table population, relational referential integrity, and data ranges, outputting validation reports."
]

for i, bullet in enumerate(bullets):
    p = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
    format_bullet(p, bullet)


# ==========================================
# SLIDE 6: Quantitative Returns & Metrics
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_slide_header(slide, "Quantitative Returns & Volatility Calculations")

left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(5.0))
ltf = left_box.text_frame
ltf.word_wrap = True

bullets = [
    "Trailing Daily NAV Returns: Calculated daily return series using: Return(t) = (NAV(t) - NAV(t-1)) / NAV(t-1). Enforced forward-filled NAVs to guarantee daily return continuity on consecutive business days.",
    "Compounded Annual Growth Rate (CAGR): Formulated the compound growth rate: CAGR = (Ending Value / Beginning Value) ^ (365.25 / Days) - 1. Calculates annualized performance across the 4.5-year history.",
    "Maximum Drawdown (MDD): Designed to measure peak-to-trough drops: Drawdown(t) = (NAV(t) - Peak(t)) / Peak(t). The MDD is the minimum observed drawdown value, serving as a worst-case risk proxy."
]
for i, bullet in enumerate(bullets):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    format_bullet(p, bullet)

# Right Column - Formulas & Callouts
formula_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.0))
rtf = formula_box.text_frame
rtf.word_wrap = True

# Add mathematical callout text
rp = rtf.paragraphs[0]
rp.text = "MATHEMATICAL FORMULAS"
rp.font.name = 'Georgia'
rp.font.size = Pt(16)
rp.font.bold = True
rp.font.color.rgb = PRIMARY_DARK
rp.space_after = Pt(20)

formulas = [
    "Daily Return Formula:\n  R_t = (NAV_t - NAV_{t-1}) / NAV_{t-1}",
    "CAGR Formula:\n  CAGR = (V_end / V_begin) ^ (365.25 / d) - 1",
    "Maximum Drawdown Formula:\n  MDD = Min((NAV_t - Peak_t) / Peak_t)"
]
for f in formulas:
    p = rtf.add_paragraph()
    p.text = f
    p.font.name = 'Consolas'
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(20)


# ==========================================
# SLIDE 7: Risk-Adjusted Ratios (CAPM)
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_slide_header(slide, "Risk-Adjusted Performance & CAPM Model")

left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(5.0))
ltf = left_box.text_frame
ltf.word_wrap = True

bullets = [
    "Sharpe Ratio: Measures risk-adjusted excess return: Sharpe = (R_annualized - R_f) / Vol_annualized. Utilizes a standardized 6.0% annual risk-free rate (R_f), annualizing daily returns and standard deviations using 252 business days.",
    "Sortino Ratio: Focuses solely on downside volatility: Sortino = (R_annualized - R_f) / Vol_downside. Downside deviation is computed strictly from negative daily returns, penalizing only negative variations.",
    "CAPM Beta (β): Measures systemic market volatility matching daily fund returns against its assigned benchmark index (e.g. NIFTY 100): Beta = Cov(R_fund, R_market) / Var(R_market).",
    "CAPM Alpha (α): Measures manager outperformance against the market: Alpha = R_fund - [R_f + Beta * (R_market - R_f)]."
]
for i, bullet in enumerate(bullets):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    format_bullet(p, bullet)

# Right Column - Formulas & Callouts
formula_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.0))
rtf = formula_box.text_frame
rtf.word_wrap = True

rp = rtf.paragraphs[0]
rp.text = "RISK MODEL EQUATIONS"
rp.font.name = 'Georgia'
rp.font.size = Pt(16)
rp.font.bold = True
rp.font.color.rgb = PRIMARY_DARK
rp.space_after = Pt(20)

formulas = [
    "Sharpe Ratio:\n  (Mean(R_d)*252 - 0.06) / (Std(R_d)*sqrt(252))",
    "Sortino Ratio:\n  (Mean(R_d)*252 - 0.06) / (Std(R_neg)*sqrt(252))",
    "CAPM Beta (β):\n  Covariance(R_fund, R_index) / Variance(R_index)",
    "CAPM Alpha (α):\n  R_fund_ann - [R_f + β * (R_index_ann - R_f)]"
]
for f in formulas:
    p = rtf.add_paragraph()
    p.text = f
    p.font.name = 'Consolas'
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(15)


# ==========================================
# SLIDE 8: Advanced Risk & Strategy Analytics
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_slide_header(slide, "Advanced Portfolio Risk & Strategy Analytics")

main_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
mtf = main_box.text_frame
mtf.word_wrap = True

bullets = [
    "Value at Risk (VaR): Calculated daily 95% VaR (5th percentile of daily return history). Represents the maximum expected loss over a single business day with a 95% confidence level.",
    "Conditional Value at Risk (CVaR): Expected loss in the worst 5% tail events (Expected Shortfall). Measures the severity of extreme downside losses (where loss exceeds VaR), capturing systemic tail-risk.",
    "Herfindahl-Hirschman Index (HHI): Formulated HHI = Sum(w_i^2) to assess sector concentration. Categorized as Diversified (HHI < 1,500), Moderately Concentrated (1,500 - 2,500), or Highly Concentrated (HHI > 2,500).",
    "Cohort Retention Decay: Categorizes investors into cohorts based on onboarding quarter (e.g. 2024Q1). Tracks active participation decay rates at 3, 6, and 12-month marks to measure customer lifecycle duration.",
    "SIP Continuation Alerts: Classifies systematic investment plans based on transaction recency. Accounts without transactions for 35+ days are flagged as At Risk to trigger proactive client outreach."
]

for i, bullet in enumerate(bullets):
    p = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
    format_bullet(p, bullet)


# ==========================================
# SLIDE 9: Key Project Metrics & Payment Mode Chart
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_slide_header(slide, "Key Project Metrics: Database Summary & Inflows")

# Left Column (Metrics)
left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(5.0))
ltf = left_box.text_frame
ltf.word_wrap = True

bullets = [
    "Total Assets Under Management: Encompasses 1,043,664.00 Crore INR (~10.4 Lakh Crore) across the 40 mutual fund schemes.",
    "Transaction Auditing: Successfully verified 32,778 individual transactions spanning 5,000 unique active investors.",
    "Cumulative Capital Inflows: Aggregated transaction value stands at 3,521,580,430.00 INR (~352 Crore INR).",
    "Digital Payment Penetration: Digital modes (UPI: 888.24M, Net Banking: 893.49M) have reached complete volume and value parity with traditional methods (Cheque: 892.22M, Mandate: 847.63M).",
    "Digitalization Trend: UPI accounts for 8,154 transactions, closely matching Net Banking (8,250 transactions) and Cheque (8,228 transactions), signaling high financial digitization."
]
for i, bullet in enumerate(bullets):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    format_bullet(p, bullet)

# Right Column (Payment Mode Chart)
insert_chart(
    slide,
    "reports/charts/investor_payment_mode.png",
    Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.3),
    "Fig 1: Inflow volume and value distribution across payment channels"
)


# ==========================================
# SLIDE 10: Mutual Fund Scorecard & NAV Growth Chart
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_slide_header(slide, "Mutual Fund Scorecard: Top & Bottom Performers")

# Left Column (Scorecard Analysis)
left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(5.0))
ltf = left_box.text_frame
ltf.word_wrap = True

bullets = [
    "Weighted Scorecard formulation: Designed composite scores utilizing Sharpe (30%), CAPM Alpha (25%), CAGR (20%), Max Drawdown (15%), and Expense Ratio (10%).",
    "Top Performing Funds: Mirae Asset Large Cap Fund (CAGR: 14.81%, Sharpe: 1.760, MDD: -11.27%), Kotak Flexicap Fund (CAGR: 15.65%, Sharpe: 1.568), and ICICI Pru Midcap Fund (CAGR: 18.08%, Sharpe: 1.391).",
    "Bottom Performing Funds: Axis Small Cap Fund (CAGR: 1.52%, Sharpe: -0.179, MDD: -51.68%), UTI Mid Cap Fund (Sharpe: -0.266), and SBI Small Cap Fund (CAGR: 2.05%, MDD: -52.57%).",
    "Performance Drivers: Mid-cap and Large-cap Equity funds dominated the scorecards due to strong bull market growth. Small-caps scored poorly due to severe maximum drawdowns (>50%)."
]
for i, bullet in enumerate(bullets):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    format_bullet(p, bullet)

# Right Column (NAV Growth Chart)
insert_chart(
    slide,
    "reports/charts/nav_growth.png",
    Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.3),
    "Fig 2: Historical NAV growth trends of top performing mutual funds"
)


# ==========================================
# SLIDE 11: Investor Demographics
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_slide_header(slide, "Investor Demographics & Geographic Insights")

# Left Column (Demographics)
left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(5.0))
ltf = left_box.text_frame
ltf.word_wrap = True

bullets = [
    "Millennial Dominance (26-35 age cohort): Represents the largest segment by value, accounting for 1,451,600,218.00 INR (41.2% total market share).",
    "Mid-Career Segment (36-45 age cohort): Represents the second largest segment, contributing 871,647,528.00 INR (24.8% market share).",
    "Gen Z Adoption (18-25 age cohort): Represents 15.1% of transactions (531.64M INR), showcasing emerging retail interest.",
    "Geographic Leaders: Punjab leads with 2,965 transactions (315.78M INR), followed closely by Tamil Nadu (315.18M INR) and Madhya Pradesh (308.31M INR).",
    "Retail Inflow Dispersion: Small ticket transaction values are highly spread across states, reflecting a broad retail investor footprint."
]
for i, bullet in enumerate(bullets):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    format_bullet(p, bullet)

# Right Column (Demographics Chart)
insert_chart(
    slide,
    "reports/charts/investor_demographics_age_gender.png",
    Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.3),
    "Fig 3: Age and Gender breakdown of transaction volume and value"
)


# ==========================================
# SLIDE 12: Portfolio & Sector Concentration
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_slide_header(slide, "Portfolio Diversification & Sector Weightings")

# Left Column
left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(5.0))
ltf = left_box.text_frame
ltf.word_wrap = True

bullets = [
    "HHI Concentration Audit: Out of 34 equity funds, 3 are classified as Diversified (HHI < 1,500), 27 as Moderately Concentrated (1,500-2,500), and 4 as Highly Concentrated (HHI > 2,500).",
    "Average Concentration Levels: The average HHI for moderately concentrated schemes lies at 2,014.96. Highly concentrated funds average 2,641.02, indicating high stock-specific risk.",
    "Key Sector Allocations: Portfolio structures lean heavily on Banking (average weight of 10.87%, representing 62,840 Crore AUM) and IT sectors (average weight of 11.39%, 38,477 Crore AUM).",
    "Concentration Risk: Highly concentrated funds require close monitoring, as their returns are tightly coupled to banking and IT sector cycles."
]
for i, bullet in enumerate(bullets):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    format_bullet(p, bullet)

# Right Column (Sector holdings treemap)
insert_chart(
    slide,
    "reports/charts/sector_holdings_treemap.png",
    Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.3),
    "Fig 4: Sector holdings treemap showing asset allocations"
)


# ==========================================
# SLIDE 13: Cohort Retention Analysis
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_slide_header(slide, "Investor Cohort Retention Analysis")

# Left Column (Cohort Decay)
left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(5.0))
ltf = left_box.text_frame
ltf.word_wrap = True

bullets = [
    "Cohort Grouping: Monitored active investor decay quarterly based on onboarding cohorts to capture structural retention patterns.",
    "2024Q1 Cohort (High Stability): Started with 3,236 investors. Exhibits strong retention: 70.92% active after 3 months, stabilizing at 69.62% active after 12 months.",
    "2024Q2 Cohort (Moderate Decay): Onboarded 971 investors. Retention declines to 63.13% active after 3 months, dropping to 48.71% active at Year 1.",
    "2024Q4 Cohort (Critical Churn): Onboarded 186 investors. Severe decay observed: only 40.32% active at 3 months, falling to 26.34% active at 6 months.",
    "Strategic Conclusion: Year-end cohorts pose high attrition risk, necessitating engagement campaigns within their first 60 days."
]
for i, bullet in enumerate(bullets):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    format_bullet(p, bullet)

# Right Column (SIP Accounts Growth Chart)
insert_chart(
    slide,
    "reports/charts/sip_accounts_active_vs_new.png",
    Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.3),
    "Fig 5: Growth of active versus newly registered systematic investment plans"
)


# ==========================================
# SLIDE 14: SIP Continuation Risk
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_slide_header(slide, "SIP Continuation & Investor Churn Risks")

# Left Column (SIP Continuation)
left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.8), Inches(5.0))
ltf = left_box.text_frame
ltf.word_wrap = True

bullets = [
    "Continuation Classification: Classified SIP accounts based on transaction recency. Active SIPs have transacted in the last 35 days; At Risk SIPs show no activity for 35+ days.",
    "Continuation Metrics: Active SIP Investors: 1,241 (26.1%). At Risk SIP Investors: 3,521 (73.9%).",
    "Churn Hazard: More than 73% of SIP investors have lapsed, indicating a high number of paused or canceled mandates. This constitutes a severe cash flow retention risk.",
    "Volatility Impacts: The high number of at-risk SIPs corresponds to periods of short-term market consolidation, causing retail investors to pause plans.",
    "Strategic Imperative: Re-engaging at-risk accounts is the most critical commercial priority to stabilize recurring capital inflows."
]
for i, bullet in enumerate(bullets):
    p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
    format_bullet(p, bullet)

# Right Column (SIP Monthly Inflows Chart)
insert_chart(
    slide,
    "reports/charts/sip_monthly_inflows.png",
    Inches(6.8), Inches(1.7), Inches(5.8), Inches(4.3),
    "Fig 6: Monthly SIP inflows (Crore INR) showing growth trajectory"
)


# ==========================================
# SLIDE 15: Strategic Recommendations
# ==========================================
slide = prs.slides.add_slide(slide_layout)
add_slide_header(slide, "Strategic Recommendations & Action Plan")

# Four quadrant-like boxes or bullet points
rec_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.2))
rtf = rec_box.text_frame
rtf.word_wrap = True

recs = [
    "Millennial Engagement & UPI Focus: Since investors aged 26-35 represent 41.2% (1,451M INR) of total transaction value, prioritize mobile-first interfaces, quick UPI payments, and seamless digital onboarding.",
    "Q4 Cohort Retention Push: Cohorts onboarded in Q4 show severe drop-offs (only 26.3% active after 6 months). Deploy targeted push notifications, educational newsletters, and reward incentives in their first 60 days.",
    "Automated SIP Resumption Prompts: To combat the 73.9% At-Risk SIP rate, implement automated triggers. When an investor skips a scheduled transaction (no contribution in 30 days), prompt them with one-click 'Resume' or 'Pause for 3 Months' actions.",
    "HHI Concentration Alerts: For investors holding 'Highly Concentrated' portfolios (HHI > 2,500), implement dashboard recommendations proposing supplementary 'Diversified' funds (HHI < 1,500) to balance holdings risk.",
    "Next Steps: Embed these metrics and automated notifications directly into the dashboard platform. Monitor monthly cohort decay to assess the effectiveness of the targeted retention campaigns."
]

for i, rec in enumerate(recs):
    p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
    format_bullet(p, rec, font_size=13)

# Ensure the output directory exists
os.makedirs("reports", exist_ok=True)

# Save presentation
output_path = "reports/capstone_project_presentation.pptx"
prs.save(output_path)
print(f"Presentation saved successfully to {output_path}")
